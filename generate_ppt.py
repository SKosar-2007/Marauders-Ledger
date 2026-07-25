#!/usr/bin/env python3
"""Generate a Harry Potter-themed demo presentation for The Marauder's Ledger."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Color Palette ===
PARCHMENT = RGBColor(0xF5, 0xE6, 0xC8)
DARK_PARCHMENT = RGBColor(0xE8, 0xD5, 0xB0)
INK = RGBColor(0x2C, 0x18, 0x10)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
BLOOD_RED = RGBColor(0xDC, 0x26, 0x26)
EMERALD = RGBColor(0x2D, 0x6A, 0x4F)
DARK_BG = RGBColor(0x1A, 0x0F, 0x0A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xD5, 0xC5)
SLIDE_BG = RGBColor(0xF8, 0xF0, 0xE0)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr is not None:
                etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val=str(alpha))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=INK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=INK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
    return txBox


def add_accent_bar(slide, left, top, width, height, color=GOLD):
    return add_shape(slide, left, top, width, height, color)


def add_card(slide, left, top, width, height, title, body_lines, title_color=INK, body_color=INK):
    add_shape(slide, left, top, width, height, WHITE)
    add_shape(slide, left, top, Inches(0.08), height, GOLD)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.3), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    add_bullet_text(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                    body_lines, font_size=11, color=body_color, spacing=Pt(3))


def add_badge(slide, left, top, text, bg_color=GOLD, text_color=DARK_BG):
    w, h = Inches(1.6), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = True
    return shape


# === Build Presentation ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank layout


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Decorative gold border
add_shape(slide, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9), DARK_BG)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(7.1), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_shape(slide, Inches(12.933), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)

# Title
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.2),
             "The Marauder's Ledger", font_size=52, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Georgia")

# Accent line
add_shape(slide, Inches(4.5), Inches(2.8), Inches(4.333), Inches(0.04), GOLD)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(0.8),
             "AI-Powered Financial Anomaly Detection", font_size=28, color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.6),
             '"I solemnly swear that I am up to no good."', font_size=20, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# Badges
add_badge(slide, Inches(3.5), Inches(5.2), "React 19", RGBColor(0x3B, 0x82, 0xF6), WHITE)
add_badge(slide, Inches(5.3), Inches(5.2), "FastAPI", EMERALD, WHITE)
add_badge(slide, Inches(7.1), Inches(5.2), "ML Ensemble", BLOOD_RED, WHITE)
add_badge(slide, Inches(8.9), Inches(5.2), "Gemini AI", AMBER, DARK_BG)

# Bottom text
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.333), Inches(0.5),
             "Hackathon Demo Presentation", font_size=14, color=RGBColor(0x99, 0x88, 0x77), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Problem & Solution
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

# Header bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "The Problem & Our Solution", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

# Problem card
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), WHITE)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6), BLOOD_RED)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.2), Inches(0.5),
             "THE PROBLEM", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4.2), [
    "Financial fraud costs billions annually",
    "Traditional rule-based systems miss novel patterns",
    "Manual review is slow and inconsistent",
    "Customers need real-time, transparent alerts",
    "Existing tools lack context and narrative explanations",
], font_size=15, color=INK, spacing=Pt(12))

# Solution card
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), WHITE)
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.6), EMERALD)
add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.5),
             "OUR SOLUTION", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(7.1), Inches(2.3), Inches(5.4), Inches(4.2), [
    "10-model ML ensemble detects complex anomalies",
    "48 engineered features for deep analysis",
    "AI-generated narratives explain each anomaly",
    "Interactive map visualizes spending patterns",
    "Voice narration for accessibility",
    "Real-time processing via Celery workers",
], font_size=15, color=INK, spacing=Pt(12))


# ============================================================
# SLIDE 3: Architecture Overview
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "System Architecture", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

# Architecture flow boxes
layers = [
    ("FRONTEND", Inches(0.5), RGBColor(0x3B, 0x82, 0xF6), [
        "React 19 + TypeScript",
        "18 Pages, 26 Components",
        "Framer Motion Animations",
        "Tailwind CSS Parchment Theme",
    ]),
    ("API LAYER", Inches(3.6), EMERALD, [
        "FastAPI + Uvicorn",
        "17+ REST Endpoints",
        "JWT Authentication",
        "CORS & Rate Limiting",
    ]),
    ("ML ENGINE", Inches(6.7), AMBER, [
        "10-Model Ensemble",
        "48 Feature Engineering",
        "Hybrid Scoring Formula",
        "Rule-Based Layer",
    ]),
    ("DATA & AI", Inches(9.8), BLOOD_RED, [
        "Actian VectorAI",
        "Gemini 2.0 Flash",
        "ElevenLabs TTS",
        "Redis + Celery",
    ]),
]

for title, left, color, items in layers:
    # Box background
    add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(5.2), WHITE)
    # Header
    add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(0.6), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.55), Inches(2.6), Inches(0.5),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Items
    add_bullet_text(slide, left + Inches(0.2), Inches(2.3), Inches(2.4), Inches(4.0),
                    items, font_size=13, color=INK, spacing=Pt(10))

# Arrows between boxes
for i in range(3):
    x = Inches(3.4 + i * 3.1)
    add_text_box(slide, x, Inches(3.5), Inches(0.4), Inches(0.5),
                 ">>>", font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom: Docker
add_shape(slide, Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.45), RGBColor(0x24, 0x96, 0xED))
add_text_box(slide, Inches(0.5), Inches(6.87), Inches(12.333), Inches(0.4),
             "Docker Compose: 5 Services  |  Backend  |  Frontend  |  Redis  |  VectorAI  |  Celery Worker",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: ML Engine
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "ML Anomaly Detection Engine", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

# Supervised models card
add_card(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(3.0),
         "SUPERVISED MODELS (6)", [
             "Random Forest",
             "Gradient Boosting",
             "XGBoost",
             "LightGBM",
             "CatBoost",
             "Extra Trees",
         ], title_color=EMERALD)

# Unsupervised models card
add_card(slide, Inches(4.8), Inches(1.4), Inches(4.0), Inches(3.0),
         "UNSUPERVISED MODELS (3)", [
             "Isolation Forest",
             "Local Outlier Factor",
             "One-Class SVM",
         ], title_color=AMBER)

# Rule-based card
add_card(slide, Inches(9.1), Inches(1.4), Inches(3.8), Inches(3.0),
         "RULE-BASED SCORING", [
             "Transaction velocity checks",
             "Amount threshold violations",
             "Category spending spikes",
             "Time-pattern anomalies",
         ], title_color=BLOOD_RED)

# Feature engineering
add_shape(slide, Inches(0.5), Inches(4.7), Inches(6.2), Inches(2.6), WHITE)
add_shape(slide, Inches(0.5), Inches(4.7), Inches(6.2), Inches(0.5), INK)
add_text_box(slide, Inches(0.8), Inches(4.75), Inches(5.6), Inches(0.4),
             "48 ENGINEERED FEATURES", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.3), Inches(5.6), Inches(1.8), [
    "35 Base Features: amount, category, merchant, hour, day, cyclical encodings",
    "3 Unsupervised Scores: IF score, LOF score, OCSVM score",
    "1 Rule Score: heuristic-based anomaly indicators",
    "Feature scaling: StandardScaler normalization",
], font_size=12, color=INK, spacing=Pt(6))

# Hybrid formula
add_shape(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(2.6), WHITE)
add_shape(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(0.5), GOLD)
add_text_box(slide, Inches(7.3), Inches(4.75), Inches(5.2), Inches(0.4),
             "HYBRID SCORING FORMULA", font_size=14, color=DARK_BG, bold=True)
add_text_box(slide, Inches(7.3), Inches(5.4), Inches(5.2), Inches(0.5),
             "final_score = 0.75 * ensemble_score + 0.25 * rule_score",
             font_size=16, color=INK, bold=True, font_name="Courier New")
add_bullet_text(slide, Inches(7.3), Inches(6.0), Inches(5.2), Inches(1.2), [
    "Ensemble: Weighted average of 9 ML models",
    "Rules: Domain-specific anomaly heuristics",
    "Threshold: Score > 0.5 flags as anomaly",
], font_size=12, color=INK, spacing=Pt(6))


# ============================================================
# SLIDE 5: Interactive Marauder's Map
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "Interactive Marauder's Map", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(3), Inches(0.04), GOLD)
add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.5),
             "SVG-based visualization with animated footprints and real-time anomaly detection",
             font_size=16, color=LIGHT_GRAY)

# Map location cards
locations = [
    ("Hogwarts", "Food & Dining", RGBColor(0x7C, 0x3A, 0xED), Inches(0.5)),
    ("Hogsmeade", "Shopping", RGBColor(0x25, 0x63, 0xEB), Inches(3.2)),
    ("Gringotts", "Bills & Utilities", RGBColor(0xD4, 0xAF, 0x37), Inches(5.9)),
    ("Diagon Alley", "Entertainment", RGBColor(0x05, 0x96, 0x69), Inches(8.6)),
    ("Platform 9 3/4", "Travel", RGBColor(0xDC, 0x26, 0x26), Inches(11.3)),
]

for name, category, color, left in locations:
    add_shape(slide, left, Inches(1.9), Inches(2.5), Inches(1.5), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.95), Inches(2.3), Inches(0.5),
                 name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.4), Inches(2.3), Inches(0.4),
                 category, font_size=12, color=RGBColor(0xFF, 0xFF, 0xCC), alignment=PP_ALIGN.CENTER)

# Feature list
features = [
    ("Animated Footprints", "Transactions appear as walking footprints that traverse the map"),
    ("Anomaly Pulsing", "Detected anomalies glow red and pulse to draw attention"),
    ("Click to Inspect", "Click any footprint to view detailed transaction information"),
    ("Category Filtering", "Filter by Moony, Wormtail, Padfoot, Prongs tabs"),
    ("Severity Indicators", "Peeves (low), Boggart (medium), Dementor (high)"),
    ("Real-time Updates", "Background Celery workers process and update the map live"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(3.7 + row * 1.7)
    add_shape(slide, left, top, Inches(3.9), Inches(1.4), RGBColor(0x2C, 0x18, 0x10))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.6), Inches(0.35),
                 title, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), Inches(3.6), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6: AI-Powered Narratives
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "AI-Powered Narratives & Voice", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

# Gemini narrative card
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.5), WHITE)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.2), Inches(0.6), RGBColor(0x42, 0x85, 0xF4))
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.6), Inches(0.5),
             "GOOGLE GEMINI 2.0 FLASH", font_size=16, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.6), Inches(0.4),
             "Natural Language Narrative Generation", font_size=18, color=INK, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(2.8), Inches(5.6), Inches(3.5), [
    "Each anomaly gets a unique, contextual narrative",
    "Harry Potter-themed storytelling style",
    "Explains what happened, why it's suspicious",
    "Typewriter animation for engaging reveal",
    "Graceful fallback to static templates",
    "Narratives stored for historical reference",
], font_size=14, color=INK, spacing=Pt(10))

# Narrative example
add_shape(slide, Inches(0.8), Inches(5.5), Inches(5.6), Inches(1.2), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(5.55), Inches(5.2), Inches(0.3),
             "Example Narrative:", font_size=11, color=AMBER, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.85), Inches(5.2), Inches(0.8),
             "\"A suspicious transaction of $847.23 at 'Dark arts Emporium' at 3:23 AM... The Boggart stirs!\"",
             font_size=11, color=INK)

# ElevenLabs card
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5), WHITE)
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.6), AMBER)
add_text_box(slide, Inches(7.3), Inches(1.55), Inches(5.2), Inches(0.5),
             "ELEVENLABS TEXT-TO-SPEECH", font_size=16, color=WHITE, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(0.4),
             "Voice Narration with Waveform", font_size=18, color=INK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.8), Inches(5.2), Inches(3.5), [
    "Narratives read aloud via TTS API",
    "Interactive audio player with waveform",
    "Play/pause controls for each anomaly",
    "Enhances accessibility",
    "Graceful fallback when API unavailable",
    "Audio cached for instant replay",
], font_size=14, color=INK, spacing=Pt(10))

# Severity classification
add_shape(slide, Inches(7.3), Inches(5.5), Inches(5.2), Inches(1.2), LIGHT_GRAY)
add_text_box(slide, Inches(7.5), Inches(5.55), Inches(4.8), Inches(0.3),
             "SEVERITY CLASSIFICATION:", font_size=11, color=INK, bold=True)

# Severity badges
add_badge(slide, Inches(7.5), Inches(5.9), "Peeves (Low)", EMERALD, WHITE)
add_badge(slide, Inches(9.2), Inches(5.9), "Boggart (Med)", AMBER, DARK_BG)
add_badge(slide, Inches(10.9), Inches(5.9), "Dementor (High)", BLOOD_RED, WHITE)


# ============================================================
# SLIDE 7: Full-Stack Features
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "Full-Stack Feature Set", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

features = [
    ("Authentication", "JWT-based auth with bcrypt\npassword hashing, CORS, and\nrate limiting", EMERALD),
    ("Background Processing", "Celery workers + Redis for\nasync CSV processing,\nanomaly detection, and\nnarrative generation", RGBColor(0x3B, 0x82, 0xF6)),
    ("18 Themed Pages", "Landing, Dashboard, Ledger,\nVault, Pensieve, Owl Post,\nProfile, Great Hall, and\n10+ more wizarding pages", RGBColor(0x7C, 0x3A, 0xED)),
    ("Dark Mode", "Full dark/light theme toggle\nwith seamless transitions\nand consistent design", RGBColor(0x64, 0x74, 0x8B)),
    ("PDF Export", "Export anomaly reports and\ntransaction data as\nformatted PDF documents", AMBER),
    ("Command Palette", "Keyboard-driven Cmd+K\ncommand interface for\nquick navigation and\nactions", INK),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.4 + row * 3.0)
    add_shape(slide, left, top, Inches(3.9), Inches(2.7), WHITE)
    add_shape(slide, left, top, Inches(3.9), Inches(0.55), color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.08), Inches(3.5), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.5), Inches(1.8),
                 desc, font_size=13, color=INK)

# Bottom stats
add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.04), GOLD)


# ============================================================
# SLIDE 8: Live Demo Flow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "Live Demo Walkthrough", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(3), Inches(0.04), GOLD)

steps = [
    ("1", "Upload", "Drag & drop a CSV file\nonto the upload zone", RGBColor(0x3B, 0x82, 0xF6)),
    ("2", "Process", "Background workers parse\n& run ML inference", EMERALD),
    ("3", "Map", "Footprints animate across\nthe Marauder's Map", AMBER),
    ("4", "Detect", "Anomalies pulse red on\nthe interactive SVG map", BLOOD_RED),
    ("5", "Narrate", "Gemini generates AI\nnarrative for each anomaly", RGBColor(0x7C, 0x3A, 0xED)),
    ("6", "Listen", "ElevenLabs reads the\nnarration via voice TTS", RGBColor(0xEC, 0x48, 0x99)),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.4 + i * 2.1)
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.7), Inches(1.6), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, left, Inches(2.5), Inches(2.0), Inches(0.4),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, left, Inches(2.95), Inches(2.0), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Additional features section
add_text_box(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             "Additional Features Explored During Demo", font_size=20, color=GOLD, bold=True)

extra_features = [
    ("MischiefList", "Historical data table with\nsearch & filter"),
    ("Vault", "Gringotts-style bank\nstatement overview"),
    ("Pensieve", "Deep spending analysis\nwith time filters"),
    ("Owl Post", "Notification center with\nread/unread states"),
    ("Profile", "Wizard's dossier with\nstats & skills"),
    ("Settings", "Admin controls &\nsystem configuration"),
]

for i, (title, desc) in enumerate(extra_features):
    col = i % 6
    left = Inches(0.3 + col * 2.15)
    add_shape(slide, left, Inches(4.8), Inches(2.0), Inches(1.8), RGBColor(0x2C, 0x18, 0x10))
    add_text_box(slide, left + Inches(0.1), Inches(4.85), Inches(1.8), Inches(0.35),
                 title, font_size=13, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(5.2), Inches(1.8), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: Performance Metrics
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             "Model Performance & Metrics", font_size=36, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)

# Big metric cards
metrics = [
    ("F1 Score", "0.873", "Test Set", EMERALD),
    ("Precision", "0.977", "Test Set", RGBColor(0x3B, 0x82, 0xF6)),
    ("CV F1", "0.8808", "Mean +/- 0.0017", AMBER),
    ("Features", "48", "Engineered", RGBColor(0x7C, 0x3A, 0xED)),
    ("Models", "10", "Ensemble", BLOOD_RED),
]

for i, (label, value, sub, color) in enumerate(metrics):
    left = Inches(0.4 + i * 2.55)
    add_shape(slide, left, Inches(1.4), Inches(2.3), Inches(2.2), WHITE)
    add_shape(slide, left, Inches(1.4), Inches(2.3), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.6), Inches(2.1), Inches(0.4),
                 label, font_size=13, color=INK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.0), Inches(2.1), Inches(0.8),
                 value, font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, left + Inches(0.1), Inches(2.85), Inches(2.1), Inches(0.4),
                 sub, font_size=11, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

# Model breakdown table
add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.333), Inches(3.3), WHITE)
add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.5), INK)
add_text_box(slide, Inches(0.8), Inches(4.05), Inches(5), Inches(0.4),
             "ENSEMBLE MODEL BREAKDOWN", font_size=14, color=WHITE, bold=True)

# Table headers
headers = ["Model", "Type", "Role", "Status"]
col_widths = [Inches(3.0), Inches(2.5), Inches(4.0), Inches(2.5)]
col_lefts = [Inches(0.7)]
for w in col_widths[:-1]:
    col_lefts.append(col_lefts[-1] + w)

for j, (header, left) in enumerate(zip(headers, col_lefts)):
    add_text_box(slide, left, Inches(4.6), col_widths[j], Inches(0.3),
                 header, font_size=11, color=INK, bold=True)

# Table rows
rows = [
    ["Random Forest, Gradient Boosting, Extra Trees", "Supervised", "Core ensemble members", "Active"],
    ["XGBoost, LightGBM, CatBoost", "Gradient Boosting", "High-performance learners", "Active"],
    ["Isolation Forest", "Unsupervised", "Anomaly detection (unsupervised)", "Active"],
    ["Local Outlier Factor (LOF)", "Unsupervised", "Density-based detection", "Active"],
    ["One-Class SVM (OCSVM)", "Unsupervised", "Boundary-based detection", "Active"],
    ["Rule-Based Scoring", "Heuristic", "Domain-specific rules", "Active"],
]

for i, row in enumerate(rows):
    top = Inches(4.9 + i * 0.37)
    if i % 2 == 0:
        add_shape(slide, Inches(0.5), top, Inches(12.333), Inches(0.37), LIGHT_GRAY)
    for j, (cell, left) in enumerate(zip(row, col_lefts)):
        color = INK
        if j == 3:
            color = EMERALD
        add_text_box(slide, left, top + Inches(0.05), col_widths[j], Inches(0.3),
                     cell, font_size=10, color=color)


# ============================================================
# SLIDE 10: Closing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Decorative gold border
add_shape(slide, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9), DARK_BG)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(7.1), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_shape(slide, Inches(12.933), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)

# Tagline
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(0.8),
             '"I solemnly swear that I am up to no good."',
             font_size=28, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# Accent line
add_shape(slide, Inches(5.0), Inches(2.8), Inches(3.333), Inches(0.04), GOLD)

# Title
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(1.0),
             "Thank You!", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Georgia")

# Summary line
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(0.6),
             "The Marauder's Ledger - AI-Powered Financial Anomaly Detection",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Stats summary
stats_text = "10 ML Models  |  48 Features  |  18 Pages  |  F1: 0.873  |  Full-Stack"
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.5),
             stats_text, font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)

# Tech stack badges
techs = ["React 19", "FastAPI", "Celery", "VectorAI", "Gemini", "ElevenLabs", "Docker"]
for i, tech in enumerate(techs):
    left = Inches(1.5 + i * 1.5)
    add_badge(slide, left, Inches(5.8), tech, RGBColor(0x33, 0x22, 0x11), GOLD)

# Q&A
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.333), Inches(0.5),
             "Questions & Discussion", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# === Save ===
output_path = "/Users/shamina/Desktop/Hackathon/demo_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
