#!/usr/bin/env python3
"""Generate a comprehensive workflow presentation and PDF for The Marauder's Ledger."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette ===
PARCHMENT = RGBColor(0xF5, 0xE6, 0xC8)
DARK_PARCHMENT = RGBColor(0xE8, 0xD5, 0xB0)
INK = RGBColor(0x2C, 0x18, 0x10)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
BLOOD_RED = RGBColor(0xDC, 0x26, 0x26)
EMERALD = RGBColor(0x2D, 0x6A, 0x4F)
DARK_BG = RGBColor(0x1A, 0x0F, 0x0A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xD5, 0xC5)
SLIDE_BG = RGBColor(0xF8, 0xF0, 0xE0)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL = RGBColor(0x05, 0x96, 0x69)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=INK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=INK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
    return txBox


def add_accent_bar(slide, left, top, width, height, color=GOLD):
    return add_shape(slide, left, top, width, height, color)


def add_card(slide, left, top, width, height, title, body_lines, title_color=INK, body_color=INK):
    add_shape(slide, left, top, width, height, WHITE)
    add_shape(slide, left, top, Inches(0.08), height, GOLD)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.3), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    add_bullet_text(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                    body_lines, font_size=11, color=body_color, spacing=Pt(3))


def add_badge(slide, left, top, text, bg_color=GOLD, text_color=DARK_BG):
    w, h = Inches(1.6), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = True
    return shape


def add_slide_header(slide, title):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                 title, font_size=36, color=GOLD, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.04), GOLD)


# ==================== BUILD PRESENTATION ====================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9), DARK_BG)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(7.1), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_shape(slide, Inches(12.933), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.2),
             "The Marauder's Ledger", font_size=52, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")
add_shape(slide, Inches(4.5), Inches(2.8), Inches(4.333), Inches(0.04), GOLD)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(0.8),
             "AI-Powered Financial Anomaly Detection", font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.6),
             '"I solemnly swear that I am up to no good."', font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)
add_badge(slide, Inches(3.5), Inches(5.2), "React 19", BLUE, WHITE)
add_badge(slide, Inches(5.3), Inches(5.2), "FastAPI", EMERALD, WHITE)
add_badge(slide, Inches(7.1), Inches(5.2), "ML Ensemble", BLOOD_RED, WHITE)
add_badge(slide, Inches(8.9), Inches(5.2), "Gemini AI", AMBER, DARK_BG)
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.333), Inches(0.5),
             "Complete Workflow Documentation", font_size=14, color=RGBColor(0x99, 0x88, 0x77), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Website Overview
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Website Overview")

add_card(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.5),
         "WHAT IS THE MARAUDER'S LEDGER?", [
             "A full-stack AI-powered web application for detecting",
             "financial anomalies in bank transaction CSV files",
             "",
             "Users upload CSV transaction data, and the system:",
             "1. Parses and validates the data",
             "2. Runs a 10-model ML ensemble for anomaly detection",
             "3. Visualizes results on an interactive SVG map",
             "4. Generates AI narratives via Google Gemini 2.0 Flash",
             "5. Reads narratives aloud via ElevenLabs TTS",
             "",
             "All themed around the Harry Potter Marauder's Map",
             "with terminology like: Dementor (high severity),",
             "Boggart (medium), Peeves (low), Mischief (fraud)",
         ], title_color=INK)

add_card(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5),
         "TECH STACK", [
             "Frontend: React 19 + TypeScript + Vite",
             "State: TanStack Query v5 + React Context",
             "Styling: Tailwind CSS v4 + Framer Motion",
             "Backend: FastAPI + Python 3.11 + Uvicorn",
             "Database: Actian VectorAI + SQLite fallback",
             "ML: scikit-learn, XGBoost, LightGBM, CatBoost",
             "AI: Google Gemini 2.0 Flash (narratives + chat)",
             "TTS: ElevenLabs API (voice narration)",
             "Async: Celery + Redis (background processing)",
             "Infra: Docker Compose (5 services)",
             "",
             "Dual frontend: Harry Potter theme (frontend/)",
             "and OmniLedger tech theme (stitch-frontend/)",
         ], title_color=INK)


# ============================================================
# SLIDE 3: User Onboarding Flow (Authentication)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 1: Authentication & Onboarding")

# Flow diagram using shapes
steps = [
    ("1", "Visit Site", "User navigates to the\napplication URL", BLUE),
    ("2", "Login/Signup", "Register with name,\nemail, password", EMERALD),
    ("3", "JWT Issued", "Backend returns JWT\ntoken (24h expiry)", AMBER),
    ("4", "Onboarding Tour", "First-time users see\na 5-step guided tour", PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * 3.2)
    add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(2.8), WHITE)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.0), Inches(1.6), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(28)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, left + Inches(0.1), Inches(2.6), Inches(2.6), Inches(0.5),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.1), Inches(2.6), Inches(0.9),
                 desc, font_size=13, color=INK, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, left + Inches(2.8), Inches(2.5), Inches(0.5), Inches(0.5),
                     ">>", font_size=20, color=GOLD, bold=True)

# Auth details card
add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.7), WHITE)
add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.45), INK)
add_text_box(slide, Inches(0.8), Inches(4.65), Inches(5), Inches(0.4),
             "AUTHENTICATION DETAILS", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.8), [
    "Registration: POST /api/auth/register (bcrypt hashing)",
    "Login: POST /api/auth/login (returns JWT + user info)",
    "Token stored in localStorage key 'marauders_token'",
    "All API calls include Authorization: Bearer <token>",
    "ProtectedRoute redirects to /login if unauthenticated",
    "401 response auto-clears session and redirects",
], font_size=13, color=INK, spacing=Pt(6))
add_bullet_text(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(1.8), [
    "Rate limiting: 30 requests/60s per IP on auth",
    "Backend verifies via get_current_user dependency",
    "User ownership checked on all endpoints (403 if mismatch)",
    "Session persists across page refreshes",
    "Logout clears token and redirects to login",
    "No refresh tokens; simple 24-hour JWT expiry",
], font_size=13, color=INK, spacing=Pt(6))


# ============================================================
# SLIDE 4: CSV Upload & Processing Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 2: CSV Upload & Processing")

steps = [
    ("1", "Upload CSV", "Drag-drop or browse\nCSV file on Landing\npage (UploadZone)", BLUE),
    ("2", "Parse & Validate", "Frontend sends to\nPOST /api/upload\nBackend parses CSV", EMERALD),
    ("3", "Celery Task", "process_upload task\ndispatched to Celery\nworker via Redis", AMBER),
    ("4", "Feature Engineering", "48 features computed\nper transaction\n(amount, hour, etc.)", PURPLE),
    ("5", "ML Inference", "10-model ensemble\nruns anomaly detection\non all transactions", BLOOD_RED),
    ("6", "Results Stored", "Anomalies inserted\ninto VectorAI/SQLite\nBatch marked complete", TEAL),
]

for i, (num, title, desc, color) in enumerate(steps):
    col, row = i % 3, i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.5 + row * 2.8)
    add_shape(slide, left, top, Inches(3.9), Inches(2.5), WHITE)
    add_shape(slide, left, top, Inches(3.9), Inches(0.06), color)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.4), top + Inches(0.2), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(22)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, left + Inches(0.1), top + Inches(0.9), Inches(3.7), Inches(0.4),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(1.3), Inches(3.7), Inches(1.0),
                 desc, font_size=12, color=INK, alignment=PP_ALIGN.CENTER)

# Polling detail
add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4),
             "Frontend polls GET /api/batches/:batch_id/progress every 1 second (up to 60s timeout) → navigates to /dashboard on completion",
             font_size=12, color=INK, alignment=PP_ALIGN.CENTER, bold=True)


# ============================================================
# SLIDE 5: ML Anomaly Detection Engine - Deep Dive
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 3: ML Anomaly Detection Engine")

# Left: Model cards
add_card(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(3.0),
         "SUPERVISED MODELS (Weighted)", [
             "XGBoost (weight: 0.25)",
             "LightGBM (weight: 0.25)",
             "Random Forest (weight: 0.15)",
             "Gradient Boosting (weight: 0.15)",
             "CatBoost (weight: 0.10)",
             "Extra Trees (weight: 0.10)",
         ], title_color=BLUE)

add_card(slide, Inches(4.8), Inches(1.4), Inches(4.0), Inches(3.0),
         "UNSUPERVISED MODELS", [
             "Isolation Forest (outlier score)",
             "Local Outlier Factor (density score)",
             "One-Class SVM (boundary score)",
             "Scores used as additional features",
             "Unsupervised scores fed into",
             "the ensemble as features",
         ], title_color=AMBER)

add_card(slide, Inches(9.1), Inches(1.4), Inches(3.8), Inches(3.0),
         "RULE-BASED SCORING", [
             "Transaction velocity checks",
             "Amount threshold violations",
             "Category spending spikes",
             "Unusual hour patterns",
             "New merchant detection",
             "Rolling average deviation",
         ], title_color=BLOOD_RED)

# Bottom: Feature engineering and formula
add_shape(slide, Inches(0.5), Inches(4.7), Inches(6.2), Inches(2.6), WHITE)
add_shape(slide, Inches(0.5), Inches(4.7), Inches(6.2), Inches(0.5), INK)
add_text_box(slide, Inches(0.8), Inches(4.75), Inches(5.6), Inches(0.4),
             "48 ENGINEERED FEATURES", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.3), Inches(5.6), Inches(1.8), [
    "35 Base Features: amount, category, merchant, hour, day,",
    "  cyclical time encoding (sin/cos), log transforms,",
    "  Z-scores, rolling statistics (mean, std, count),",
    "  interaction features (category×amount, hour×day)",
    "3 Unsupervised Scores: IF score, LOF score, OCSVM score",
    "1 Rule Score: heuristic-based anomaly indicators",
    "Feature scaling: StandardScaler normalization",
], font_size=12, color=INK, spacing=Pt(4))

add_shape(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(2.6), WHITE)
add_shape(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(0.5), GOLD)
add_text_box(slide, Inches(7.3), Inches(4.75), Inches(5.2), Inches(0.4),
             "HYBRID SCORING FORMULA", font_size=14, color=DARK_BG, bold=True)
add_text_box(slide, Inches(7.3), Inches(5.4), Inches(5.2), Inches(0.5),
             "final_score = 0.75 * ensemble + 0.25 * rules",
             font_size=16, color=INK, bold=True, font_name="Courier New")
add_bullet_text(slide, Inches(7.3), Inches(5.9), Inches(5.2), Inches(1.2), [
    "Ensemble: Weighted average of 9 ML models",
    "Rules: Domain-specific anomaly heuristics",
    "Threshold: Score > 0.5 flagged as anomaly",
], font_size=12, color=INK, spacing=Pt(6))


# ============================================================
# SLIDE 6: Dashboard & Marauder's Map Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "Workflow Phase 4: Dashboard & Interactive Map", font_size=32, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(3), Inches(0.04), GOLD)
add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.5),
             "SVG-based Marauder's Map with animated footprints and real-time anomaly detection",
             font_size=16, color=LIGHT_GRAY)

# Map location cards
locations = [
    ("Hogwarts", "Food & Dining", PURPLE, Inches(0.5)),
    ("Hogsmeade", "Shopping", BLUE, Inches(3.2)),
    ("Gringotts", "Bills & Utilities", GOLD, Inches(5.9)),
    ("Diagon Alley", "Entertainment", TEAL, Inches(8.6)),
    ("Platform 9 3/4", "Travel", BLOOD_RED, Inches(11.3)),
]
for name, category, color, left in locations:
    add_shape(slide, left, Inches(1.9), Inches(2.5), Inches(1.5), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.95), Inches(2.3), Inches(0.5),
                 name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.4), Inches(2.3), Inches(0.4),
                 category, font_size=12, color=RGBColor(0xFF, 0xFF, 0xCC), alignment=PP_ALIGN.CENTER)

# Feature descriptions
features = [
    ("Animated Footprints", "Transactions appear as walking footprints\ntraversing the map between locations"),
    ("Anomaly Pulsing", "Detected anomalies glow red and pulse\nto draw immediate attention"),
    ("Severity Indicators", "Peeves (green/low), Boggart (amber/medium),\nDementor (red/high severity)"),
    ("Category Filtering", "Filter tabs: Moony (Food), Wormtail (Shopping),\nPadfoot (Bills), Prongs (All)"),
    ("Spending Chart", "Bottom panel shows Recharts line chart\nof spending trends over time"),
    ("Stats Bar", "Top bar: total anomalies, high/medium/low\ncounts, total financial impact"),
]
for i, (title, desc) in enumerate(features):
    col, row = i % 3, i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(3.7 + row * 1.7)
    add_shape(slide, left, top, Inches(3.9), Inches(1.4), RGBColor(0x2C, 0x18, 0x10))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.6), Inches(0.35),
                 title, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), Inches(3.6), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 7: Anomaly Investigation Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 5: Anomaly Investigation")

# Flow steps
steps = [
    ("1", "Click Anomaly", "Click map marker or\nanomaly card to\nnavigate to detail\n/anomaly/:id", BLUE),
    ("2", "Score Gauges", "Three circular gauges:\n- ML Model Score\n- Rule Score\n- Final Mischief Score", EMERALD),
    ("3", "AI Narrative", "Gemini generates HP-\nthemed narrative\nwith typewriter\neffect animation", PURPLE),
    ("4", "Audio Narration", "ElevenLabs TTS reads\nnarrative aloud via\nwaveform audio player\n(play/pause)", AMBER),
    ("5", "Review & Act", "Mark as 'Valid' or\n'Mischief'. Export\nHTML report. Share\nvia clipboard link.", BLOOD_RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.2 + i * 2.6)
    add_shape(slide, left, Inches(1.5), Inches(2.4), Inches(3.0), WHITE)
    add_shape(slide, left, Inches(1.5), Inches(2.4), Inches(0.06), color)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(1.6), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, left + Inches(0.1), Inches(2.5), Inches(2.2), Inches(0.4),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.9), Inches(2.2), Inches(1.2),
                 desc, font_size=11, color=INK, alignment=PP_ALIGN.CENTER)

# Bottom: Detail features
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.45), INK)
add_text_box(slide, Inches(0.5), Inches(4.85), Inches(12.333), Inches(0.4),
             "ANOMALY DETAIL PAGE FEATURES", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

features = [
    ("Triggered Rules", "Tags showing which rules the\nanomaly violated (e.g., 'Unusual\nHour', 'Amount Spike')"),
    ("Related Transactions", "Nearby transactions for context\nand comparison within the\nsame time window"),
    ("Status Management", "Toggle between 'Valid' and\n'Mischief' status to provide\nfeedback for model improvement"),
    ("Report Export", "Export anomaly details as a\nformatted HTML document for\nrecord-keeping or sharing"),
]
for i, (title, desc) in enumerate(features):
    left = Inches(0.3 + i * 3.3)
    add_shape(slide, left, Inches(5.4), Inches(3.1), Inches(1.8), WHITE)
    add_shape(slide, left, Inches(5.4), Inches(3.1), Inches(0.06), GOLD)
    add_text_box(slide, left + Inches(0.15), Inches(5.5), Inches(2.8), Inches(0.35),
                 title, font_size=13, color=INK, bold=True)
    add_text_box(slide, left + Inches(0.15), Inches(5.9), Inches(2.8), Inches(1.0),
                 desc, font_size=11, color=INK)


# ============================================================
# SLIDE 8: AI Narratives & Voice Chat Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 6: AI Narratives & Voice Chat")

# Gemini narrative flow
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(3.0), WHITE)
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.5), BLUE)
add_text_box(slide, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4),
             "NARRATIVE GENERATION FLOW", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(2.1), Inches(5.6), Inches(2.0), [
    "1. User navigates to /anomaly/:id detail page",
    "2. Frontend calls GET /api/narratives/:id",
    "3. Backend checks cache for existing narrative",
    "4. If not cached: calls Google Gemini 2.0 Flash API",
    "5. Gemini generates HP-themed narrative explaining the anomaly",
    "6. Narrative stored in DB, returned to frontend",
    "7. Frontend displays with typewriter animation (NarrativeCard)",
    "8. Graceful fallback: static template if Gemini API unavailable",
], font_size=13, color=INK, spacing=Pt(4))

# ElevenLabs TTS flow
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(3.0), WHITE)
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5), AMBER)
add_text_box(slide, Inches(7.3), Inches(1.45), Inches(5.2), Inches(0.4),
             "TTS AUDIO GENERATION FLOW", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(2.0), [
    "1. User clicks play on AudioPlayer widget",
    "2. Frontend calls GET /api/narratives/:id/audio",
    "3. Backend checks cache for existing MP3 audio",
    "4. If not cached: calls ElevenLabs TTS API (voice: Roger)",
    "5. Audio streamed back as MP3 (Content-Type: audio/mpeg)",
    "6. AudioPlayer shows waveform controls (play/pause)",
    "7. Sequential request locking per anomaly ID",
    "8. Graceful fallback: narrative displayed as text only",
], font_size=13, color=INK, spacing=Pt(4))

# Voice Chat
add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.7), WHITE)
add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.45), PURPLE)
add_text_box(slide, Inches(0.8), Inches(4.65), Inches(5), Inches(0.4),
             "VOICE CHAT ASSISTANT", font_size=14, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.8), [
    "Persistent FAB (floating action button) in bottom-right",
    "Supports voice input via Web Speech API + text input",
    "Context sent: anomaly_id, batch_id for relevant responses",
    "POST /api/chat/message → Gemini generates response",
    "Response: X-Chat-Response header + audio/mpeg stream",
    "ElevenLabs reads Marauder's Map persona responses aloud",
], font_size=13, color=INK, spacing=Pt(6))
add_bullet_text(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(1.8), [
    "Ask questions like: 'What anomalies were detected?'",
    "'Explain this transaction' or 'Summarize my spending'",
    "Context-aware: responds based on current page context",
    "Fallback to text-only response if speech not supported",
    "Chat history maintained during session",
    "HP-themed persona: 'Mischief Managed!' responses",
], font_size=13, color=INK, spacing=Pt(6))


# ============================================================
# SLIDE 9: Supporting Pages Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 7: Supporting Pages & Navigation")

pages = [
    ("/ledger (MischiefList)", "Searchable, filterable, paginated\nhistorical anomaly table. Shows all\ndetected anomalies with severity\nbadges, timestamps, and amounts.", BLUE),
    ("/vault (Gringotts Vault)", "Spending overview by category with\nprogress rings. Total balance,\ncategory breakdown, recent\nanomalies in the Vault view.", EMERALD),
    ("/pensieve (Deep Analysis)", "Spending trend line charts (Recharts)\nwith time range selectors. Risk\nbreakdown pie chart, spending\ncategories with progress bars.", AMBER),
    ("/owl-post (Notifications)", "Notification center with read/unread\nanomaly alerts. Categorized by\nseverity (Dementor, Boggart,\nPeeves) with timestamps.", PURPLE),
    ("/profile (Wizard's Dossier)", "User stats: total transactions,\nanomalies found, investigations.\nSpending by category rings and\nrecent cases list.", BLOOD_RED),
    ("/admin (Settings)", "System overview, theme toggle\n(dark/light), parchment aging\neffect, mischief taxonomy editor,\nAPI status dashboard.", TEAL),
]

for i, (page, desc, color) in enumerate(pages):
    col, row = i % 3, i // 3
    left = Inches(0.3 + col * 4.3)
    top = Inches(1.4 + row * 3.0)
    add_shape(slide, left, top, Inches(4.1), Inches(2.7), WHITE)
    add_shape(slide, left, top, Inches(4.1), Inches(0.06), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(3.8), Inches(0.4),
                 page, font_size=14, color=color, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.7), Inches(3.8), Inches(1.8),
                 desc, font_size=13, color=INK)


# ============================================================
# SLIDE 10: Complete Data Flow Architecture
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "End-to-End Data Flow Architecture")

# Vertical flow diagram
layers_data = [
    ("USER BROWSER (React 19)", BLUE, [
        "BrowserRouter → AuthProvider → QueryClientProvider → AppProvider",
        "TanStack Query: useQuery + useMutation for all API data",
        "React Context: Auth (JWT), App (batchId), Toast (notifications)",
        "Components: SidebarNav, UploadZone, MaraudersMap, VoiceChatWidget",
    ]),
    ("API GATEWAY (FastAPI + Uvicorn)", EMERALD, [
        "17+ REST endpoints: auth, upload, anomalies, narratives, chat",
        "JWT verification via get_current_user dependency",
        "CORS middleware, rate limiting (30 req/60s per IP)",
        "Async handlers with HTTPX for external API calls",
    ]),
    ("BACKGROUND PROCESSING (Celery + Redis)", AMBER, [
        "process_upload: parse CSV → feature engineering → ML inference",
        "detect_anomalies: ensemble prediction + rule scoring",
        "Frontend polls GET /api/batches/:id/progress every 1s",
        "Synchronous fallback if Celery/Redis unavailable",
    ]),
    ("MACHINE LEARNING (scikit-learn + XGBoost + LightGBM + CatBoost)", BLOOD_RED, [
        "10 models: 6 supervised + 3 unsupervised + 1 rule engine",
        "48 features per transaction → scaled → ensemble prediction",
        "Hybrid score: 0.75 * ML ensemble + 0.25 * rule_score",
        "Score > 0.5 → flagged as anomaly with severity level",
    ]),
    ("DATA LAYER (Actian VectorAI + SQLite)", PURPLE, [
        "VectorAI: primary vector database (hybrid search)",
        "SQLite: automatic fallback + user auth storage",
        "Named vectors: numerical + semantic embedding modes",
        "Auto-detection of VectorAI server capabilities",
    ]),
    ("EXTERNAL AI SERVICES (Gemini + ElevenLabs)", TEAL, [
        "Gemini 2.0 Flash: narrative gen + voice chat responses",
        "ElevenLabs: TTS narration (voice: Roger, streaming MP3)",
        "Thread-safe locking per anomaly (prevents duplicate gen)",
        "Graceful degradation if API keys are missing",
    ]),
]

for i, (title, color, items) in enumerate(layers_data):
    top = Inches(1.2 + i * 1.05)
    add_shape(slide, Inches(0.3), top, Inches(2.2), Inches(0.95), color)
    add_text_box(slide, Inches(0.35), top + Inches(0.1), Inches(2.1), Inches(0.75),
                 title, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(2.6), top, Inches(10.4), Inches(0.95), WHITE)
    add_shape(slide, Inches(2.6), top, Inches(0.06), Inches(0.95), color)
    for j, item in enumerate(items):
        add_text_box(slide, Inches(2.8), top + Inches(0.05 + j * 0.22), Inches(10.0), Inches(0.22),
                     item, font_size=10, color=INK if j < 3 else DARK_BG, bold=(j == 0))


# ============================================================
# SLIDE 11: Performance Metrics & Model Accuracy
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Performance Metrics & Validation")

metrics = [
    ("F1 Score", "0.873", "Test Set", EMERALD),
    ("Precision", "0.977", "Test Set", BLUE),
    ("Recall", "0.897", "Test Set", AMBER),
    ("CV F1", "0.8808", "Mean +/- 0.0017", PURPLE),
    ("Features", "48", "Engineered", BLOOD_RED),
    ("Models", "10", "Ensemble Size", TEAL),
]

for i, (label, value, sub, color) in enumerate(metrics):
    left = Inches(0.3 + i * 2.15)
    add_shape(slide, left, Inches(1.4), Inches(2.0), Inches(2.0), WHITE)
    add_shape(slide, left, Inches(1.4), Inches(2.0), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.6), Inches(1.8), Inches(0.3),
                 label, font_size=12, color=INK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(1.9), Inches(1.8), Inches(0.7),
                 value, font_size=40, color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, left + Inches(0.1), Inches(2.65), Inches(1.8), Inches(0.3),
                 sub, font_size=10, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

# Ensemble breakdown table
add_shape(slide, Inches(0.5), Inches(3.7), Inches(12.333), Inches(3.6), WHITE)
add_shape(slide, Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.45), INK)
add_text_box(slide, Inches(0.8), Inches(3.75), Inches(5), Inches(0.4),
             "ENSEMBLE MODEL BREAKDOWN", font_size=14, color=WHITE, bold=True)

headers = ["Model", "Type", "Role", "Weight", "Status"]
col_widths = [Inches(3.0), Inches(2.0), Inches(3.5), Inches(1.5), Inches(2.0)]
col_lefts = [Inches(0.7)]
for w in col_widths[:-1]:
    col_lefts.append(col_lefts[-1] + w)

for j, (header, left) in enumerate(zip(headers, col_lefts)):
    add_text_box(slide, left, Inches(4.25), col_widths[j], Inches(0.3),
                 header, font_size=11, color=INK, bold=True)

rows = [
    ["Random Forest", "Supervised", "Core ensemble", "0.15", "Active"],
    ["Gradient Boosting", "Supervised", "Core ensemble", "0.15", "Active"],
    ["XGBoost", "Gradient Boosting", "High-perf learner", "0.25", "Active"],
    ["LightGBM", "Gradient Boosting", "High-perf learner", "0.25", "Active"],
    ["CatBoost", "Gradient Boosting", "Supplementary", "0.10", "Active"],
    ["Extra Trees", "Supervised", "Supplementary", "0.10", "Active"],
    ["Isolation Forest", "Unsupervised", "Outlier detection", "Feature", "Active"],
    ["Local Outlier Factor", "Unsupervised", "Density detection", "Feature", "Active"],
    ["One-Class SVM", "Unsupervised", "Boundary detection", "Feature", "Active"],
    ["Rule-Based Scoring", "Heuristic", "Domain rules", "0.25 weight", "Active"],
]

for i, row in enumerate(rows):
    top = Inches(4.55 + i * 0.26)
    if i % 2 == 0:
        add_shape(slide, Inches(0.5), top, Inches(12.333), Inches(0.26), LIGHT_GRAY)
    for j, (cell, left) in enumerate(zip(row, col_lefts)):
        c = INK
        if j == 4:
            c = EMERALD
        add_text_box(slide, left, top + Inches(0.02), col_widths[j], Inches(0.22),
                     cell, font_size=9, color=c)


# ============================================================
# SLIDE 12: Voice Chat & Command Palette Workflow
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Workflow Phase 8: Power User Features")

# Voice Chat
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(2.8), WHITE)
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.5), PURPLE)
add_text_box(slide, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4),
             "VOICE CHAT WIDGET", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(2.1), Inches(5.6), Inches(1.8), [
    "FAB widget in bottom-right corner of every page",
    "Voice input via Web Speech API (SpeechRecognition)",
    "Text input fallback for environments without mic",
    "Sends POST /api/chat/message with page context",
    "Gemini responds in Marauder's Map persona",
    "ElevenLabs streams audio response back",
    "Context includes anomaly_id, batch_id for relevance",
    "Chat history preserved during session",
], font_size=12, color=INK, spacing=Pt(4))

# Command Palette
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.8), WHITE)
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5), AMBER)
add_text_box(slide, Inches(7.3), Inches(1.45), Inches(5.2), Inches(0.4),
             "COMMAND PALETTE (Cmd+K)", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(1.8), [
    "Press Cmd+K (or Ctrl+K) to open command palette",
    "Quick navigation to any page in the application",
    "Search transactions, anomalies, and investigations",
    "Toggle dark mode directly from command palette",
    "Access keyboard shortcuts reference",
    "Filter anomalies by severity or category",
    "Context-aware commands based on current page",
    "Fuzzy search for fast access",
], font_size=12, color=INK, spacing=Pt(4))

# Additional features
add_shape(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.8), WHITE)
add_shape(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(0.5), EMERALD)
add_text_box(slide, Inches(0.8), Inches(4.55), Inches(5.2), Inches(0.4),
             "DARK MODE & THEMING", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.2), Inches(5.2), Inches(1.8), [
    "Full dark/light theme toggle (DarkModeToggle)",
    "Persistent theme preference in localStorage",
    "Seamless transitions via Tailwind dark mode",
    "Parchment aging effect in Harry Potter theme",
    "Consistent design across all 18+ pages",
    "Tailwind CSS v4 with custom theme configuration",
], font_size=12, color=INK, spacing=Pt(4))

add_shape(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.8), WHITE)
add_shape(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(0.5), BLOOD_RED)
add_text_box(slide, Inches(7.1), Inches(4.55), Inches(5.4), Inches(0.4),
             "PDF EXPORT & SHARING", font_size=16, color=WHITE, bold=True)
add_bullet_text(slide, Inches(7.1), Inches(5.2), Inches(5.4), Inches(1.8), [
    "Export anomaly reports as formatted HTML/PDF",
    "Share anomaly details via clipboard link",
    "PDF includes: score gauges, narrative, rules",
    "Helps with record-keeping and reporting",
    "Export function in AnomalyDetail page",
    "pdfExport.ts service handles HTML generation",
], font_size=12, color=INK, spacing=Pt(4))


# ============================================================
# SLIDE 13: Container Architecture & Deployment
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Deployment Architecture (Docker Compose)")

# Docker services visualization
services = [
    ("Frontend", "React 19 + Nginx\nPort :3000", BLUE, Inches(0.5)),
    ("Backend", "FastAPI + Uvicorn\nPort :8000", EMERALD, Inches(3.0)),
    ("Redis", "Message Broker\nPort :6379", BLOOD_RED, Inches(5.5)),
    ("Celery Worker", "Async Tasks\n(no port)", AMBER, Inches(8.0)),
    ("VectorAI", "Vector Database\nPort :8888", PURPLE, Inches(10.5)),
]

for name, desc, color, left in services:
    add_shape(slide, left, Inches(1.5), Inches(2.3), Inches(2.5), WHITE)
    add_shape(slide, left, Inches(1.5), Inches(2.3), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.7), Inches(2.1), Inches(0.4),
                 name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.2), Inches(2.1), Inches(0.9),
                 desc, font_size=12, color=INK, alignment=PP_ALIGN.CENTER)

# Connection arrows
arrows_data = [
    ("User → Frontend", Inches(0.5), Inches(4.3), BLUE),
    ("Frontend → Backend (REST API)", Inches(0.5), Inches(4.7), EMERALD),
    ("Backend → Redis (Celery Broker)", Inches(0.5), Inches(5.1), BLOOD_RED),
    ("Celery → VectorAI (Results DB)", Inches(0.5), Inches(5.5), PURPLE),
    ("Backend → Gemini/ElevenLabs (External)", Inches(0.5), Inches(5.9), AMBER),
]

for title, left, top, color in arrows_data:
    add_shape(slide, left, top, Inches(12.333), Inches(0.35), color)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.03), Inches(12), Inches(0.3),
                 title, font_size=12, color=WHITE, bold=True)

# Bottom: startup command
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8), INK)
add_text_box(slide, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.7),
             "docker-compose up --build                    # Starts all 5 services\nBackend: http://localhost:8000  |  Frontend: http://localhost:3000  |  API Docs: http://localhost:8000/docs",
             font_size=14, color=WHITE, bold=False, font_name="Courier New")


# ============================================================
# SLIDE 14: Stitch Frontend (OmniLedger) Overview
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SLIDE_BG)
add_slide_header(slide, "Bonus: OmniLedger (Stitch Frontend)")

add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.333), Inches(1.5), WHITE)
add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.06), BLUE)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(0.4),
             "ALTERNATIVE TECH/CYBERPUNK THEMED FRONTEND", font_size=18, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.8), Inches(0.8),
             "The project includes a second frontend at stitch-frontend/ with an industrial/tech aesthetic. "
             "It connects to the same FastAPI backend and offers 18+ pages including Fleet Management, "
             "Access Control, Global Feed, Integrations, Performance, and System Health monitoring. "
             "Both frontends share the same core anomaly detection engine and can be run simultaneously.",
             font_size=13, color=INK)

# Page comparison
pages_hp = ["Landing", "Dashboard", "AnomalyDetail", "MischiefList", "Vault", "Pensieve", "OwlPost", "Profile", "Admin"]
pages_on = ["Landing", "Dashboard", "Anomalies", "AnomalyDetail", "Transactions", "Vault", "Analysis", "Activity",
            "Messaging", "Settings", "System", "Fleet", "AccessLogs", "AccessControl", "GlobalFeed", "Integrations", "Performance", "Reporting"]

add_card(slide, Inches(0.5), Inches(3.2), Inches(5.8), Inches(4.0),
         "HARRY POTTER THEME (frontend/)", [
             "9 main pages with HP-themed naming",
             "Landing: CSV upload with magical UI",
             "Dashboard: SVG Marauder's Map",
             "Anomaly: score gauges + narrative",
             "Ledger: MischiefList data table",
             "Vault: Gringotts bank overview",
             "Pensieve: deep spending analysis",
             "Owl Post: magical notification center",
             "Profile: Wizard's dossier",
             "Admin: system configuration",
         ], title_color=PURPLE)

add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(4.0),
         "OMNILEDGER TECH THEME (stitch-frontend/)", [
             "18+ pages with industrial naming",
             "Dashboard: cluster map of spending",
             "Anomalies: full list with bulk actions",
             "Transactions: complete transaction log",
             "Fleet Management: spending category groups",
             "Access Control: permission management",
             "Global Feed: real-time activity stream",
             "Integrations: third-party connections",
             "Performance: system performance metrics",
             "Reporting: detailed financial reports",
         ], title_color=BLUE)


# ============================================================
# SLIDE 15: Complete Workflow Summary
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "Complete End-to-End Workflow Summary", font_size=34, color=GOLD, bold=True, font_name="Georgia")
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(3), Inches(0.04), GOLD)

# Summary steps
summary_steps = [
    ("1. AUTHENTICATION", "User signs up or logs in.\nJWT stored in localStorage.\nProtected routes guard pages.", BLUE),
    ("2. CSV UPLOAD", "Drag-drop CSV on Landing page.\nFile sent to backend. Celery\nworker processes asynchronously.", EMERALD),
    ("3. ML DETECTION", "48 features engineered per txn.\n10-model ensemble predicts.\nHybrid scoring (ML + rules).", BLOOD_RED),
    ("4. MAP DISPLAY", "Results shown on Marauder's Map.\nFootprints animate, anomalies\npulse with severity colors.", AMBER),
    ("5. INVESTIGATION", "Click anomaly for detail view.\nScore gauges, AI narrative,\nTTS audio narration.", PURPLE),
    ("6. ACTIONS", "Mark valid/mischief. Export\nreport. Share links. Voice\nchat for natural queries.", TEAL),
]

for i, (title, desc, color) in enumerate(summary_steps):
    col, row = i % 3, i // 3
    left = Inches(0.4 + col * 4.3)
    top = Inches(1.4 + row * 3.0)
    add_shape(slide, left, top, Inches(4.0), Inches(2.7), RGBColor(0x2C, 0x18, 0x10))
    add_shape(slide, left, top, Inches(4.0), Inches(0.06), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(3.7), Inches(0.4),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.7), Inches(3.7), Inches(1.8),
                 desc, font_size=13, color=LIGHT_GRAY)

# Bottom tagline
add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             '"Mischief Managed." — The Marauder\'s Ledger v1.0',
             font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 16: Closing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9), DARK_BG)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(7.1), Inches(12.633), Inches(0.05), GOLD)
add_shape(slide, Inches(0.35), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_shape(slide, Inches(12.933), Inches(0.35), Inches(0.05), Inches(6.9), GOLD)
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(0.8),
             '"I solemnly swear that I am up to no good."',
             font_size=28, color=GOLD, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.0), Inches(2.8), Inches(3.333), Inches(0.04), GOLD)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(1.0),
             "Thank You!", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(0.6),
             "The Marauder's Ledger - AI-Powered Financial Anomaly Detection",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
stats_text = "10 ML Models  |  48 Features  |  18+ Pages  |  F1: 0.873  |  Full-Stack React + FastAPI + Docker"
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.5),
             stats_text, font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)
techs = ["React 19", "FastAPI", "Celery", "VectorAI", "Gemini", "ElevenLabs", "Docker"]
for i, tech in enumerate(techs):
    add_badge(slide, Inches(1.5 + i * 1.5), Inches(5.8), tech, RGBColor(0x33, 0x22, 0x11), GOLD)
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.333), Inches(0.5),
             "Questions & Discussion", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# === Save PPTX ===
pptx_path = "/home/saikat/Marauders-Ledger/Marauders_Ledger_Workflow_Presentation.pptx"
prs.save(pptx_path)
print(f"PPTX saved: {pptx_path}")
print(f"Total slides: {len(prs.slides)}")


# ============================================================
# GENERATE PDF VERSION
# ============================================================
print("\nGenerating PDF version...")

from weasyprint import HTML

pdf_html = """
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<style>
  @page {
    size: A4 landscape;
    margin: 0.5in;
    @top-center {
      content: "The Marauder's Ledger - Workflow Documentation";
      font-family: Georgia, serif;
      color: #D4AF37;
      font-size: 9pt;
    }
    @bottom-center {
      content: "Page " counter(page) " of " counter(pages);
      font-family: Arial, sans-serif;
      color: #666;
      font-size: 8pt;
    }
  }
  * { box-sizing: border-box; }
  body {
    font-family: 'Segoe UI', Arial, sans-serif;
    color: #2C1810;
    background: #F5F0E6;
    margin: 0;
    padding: 0;
  }
  .page {
    width: 100%;
    min-height: 10in;
    padding: 0.3in;
    margin-bottom: 0.2in;
    background: #F8F0E0;
    border: 1px solid #E8D5B0;
    border-radius: 4px;
    page-break-after: always;
    position: relative;
  }
  .page:last-child { page-break-after: auto; }

  h1 {
    font-family: Georgia, serif;
    color: #D4AF37;
    font-size: 28pt;
    margin: 0 0 5px 0;
    border-bottom: 3px solid #D4AF37;
    padding-bottom: 8px;
  }
  h2 {
    font-family: Georgia, serif;
    color: #2C1810;
    font-size: 18pt;
    margin: 15px 0 8px 0;
    padding-left: 10px;
    border-left: 5px solid #D4AF37;
  }
  h3 {
    font-family: Georgia, serif;
    color: #2C1810;
    font-size: 14pt;
    margin: 10px 0 5px 0;
  }
  h4 {
    color: #2C1810;
    font-size: 12pt;
    margin: 8px 0 4px 0;
  }
  p { margin: 4px 0; font-size: 10pt; line-height: 1.4; }
  ul { margin: 4px 0; padding-left: 20px; }
  li { font-size: 10pt; line-height: 1.3; margin-bottom: 2px; }
  .badge {
    display: inline-block;
    padding: 2px 10px;
    border-radius: 3px;
    font-size: 9pt;
    font-weight: bold;
    color: white;
    margin: 2px;
  }
  .badge-gold { background: #D4AF37; color: #1A0F0A; }
  .badge-red { background: #DC2626; }
  .badge-green { background: #2D6A4F; }
  .badge-blue { background: #3B82F6; }
  .badge-amber { background: #D97706; }
  .badge-purple { background: #7C3AED; }
  .badge-teal { background: #059669; }

  .grid { display: flex; flex-wrap: wrap; gap: 12px; }
  .grid-item {
    flex: 1;
    min-width: 250px;
    background: white;
    border: 1px solid #E8D5B0;
    border-radius: 4px;
    padding: 12px;
  }
  .grid-item h4 {
    margin-top: 0;
    border-bottom: 2px solid #D4AF37;
    padding-bottom: 4px;
  }

  .flow-row { display: flex; align-items: stretch; gap: 8px; margin: 10px 0; }
  .flow-step {
    flex: 1;
    background: white;
    border: 1px solid #E8D5B0;
    border-radius: 4px;
    padding: 10px;
    text-align: center;
  }
  .flow-num {
    display: inline-block;
    width: 30px; height: 30px;
    line-height: 30px;
    border-radius: 50%;
    color: white;
    font-weight: bold;
    font-size: 14pt;
    margin-bottom: 5px;
  }
  .flow-arrow {
    font-size: 20pt;
    color: #D4AF37;
    display: flex;
    align-items: center;
    font-weight: bold;
  }

  table {
    width: 100%;
    border-collapse: collapse;
    margin: 8px 0;
    font-size: 9pt;
  }
  th {
    background: #2C1810;
    color: white;
    padding: 6px 8px;
    text-align: left;
    font-weight: bold;
  }
  td {
    padding: 4px 8px;
    border-bottom: 1px solid #E8D5B0;
  }
  tr:nth-child(even) td { background: #F5F0E6; }

  .cover {
    text-align: center;
    padding-top: 2in;
  }
  .cover h1 { font-size: 36pt; border-bottom: none; }
  .cover .subtitle { font-size: 20pt; color: #666; margin: 15px 0; }
  .cover .tagline { font-size: 16pt; color: #D4AF37; font-style: italic; margin: 20px 0; }
  .cover .meta { font-size: 11pt; color: #998877; margin-top: 40px; }

  .section-header {
    background: #1A0F0A;
    color: #D4AF37;
    padding: 8px 15px;
    border-radius: 4px;
    font-family: Georgia, serif;
    font-size: 16pt;
    margin: 15px 0 10px 0;
  }

  .two-col { display: flex; gap: 12px; }
  .two-col > div { flex: 1; }

  .tech-row { display: flex; flex-wrap: wrap; gap: 6px; margin: 8px 0; }
  .tech-tag {
    background: #2C1810;
    color: #D4AF37;
    padding: 3px 10px;
    border-radius: 3px;
    font-size: 9pt;
    font-family: 'Courier New', monospace;
  }
</style>
</head>
<body>

<div class="page">
  <div class="cover">
    <h1>The Marauder's Ledger</h1>
    <div class="subtitle">AI-Powered Financial Anomaly Detection</div>
    <div class="tagline">"I solemnly swear that I am up to no good."</div>
    <div style="margin: 30px 0;">
      <span class="badge badge-blue">React 19</span>
      <span class="badge badge-green">FastAPI</span>
      <span class="badge badge-red">ML Ensemble</span>
      <span class="badge badge-amber">Gemini AI</span>
      <span class="badge badge-purple">VectorAI</span>
      <span class="badge badge-teal">Docker</span>
    </div>
    <div class="meta">Complete Workflow Documentation</div>
  </div>
</div>

<div class="page">
  <h1>1. Project Overview</h1>
  <p>The Marauder's Ledger (also branded as <strong>OmniLedger</strong>) is a full-stack AI-powered web application that detects financial anomalies in bank transaction CSV files. Users upload their transaction data, and the system runs a sophisticated 10-model machine learning ensemble to identify suspicious patterns, visualizes them on an interactive Marauder's Map-themed SVG interface, generates AI-powered narratives via Google Gemini, and reads them aloud via ElevenLabs text-to-speech.</p>

  <h2>Technology Stack</h2>
  <div class="two-col">
    <div>
      <h4>Frontend</h4>
      <ul>
        <li>React 19 + TypeScript 6.0</li>
        <li>Vite 8 (build tool)</li>
        <li>React Router v7 (routing)</li>
        <li>TanStack Query v5 (server state)</li>
        <li>Axios (HTTP client)</li>
        <li>Tailwind CSS v4 (styling)</li>
        <li>Framer Motion (animations)</li>
        <li>Recharts (charts)</li>
        <li>react-dropzone (CSV upload)</li>
      </ul>
    </div>
    <div>
      <h4>Backend & Infrastructure</h4>
      <ul>
        <li>Python 3.11 + FastAPI + Uvicorn</li>
        <li>Actian VectorAI + SQLite (databases)</li>
        <li>Celery + Redis (background tasks)</li>
        <li>scikit-learn, XGBoost, LightGBM, CatBoost</li>
        <li>Google Gemini 2.0 Flash (AI narrative)</li>
        <li>ElevenLabs API (text-to-speech)</li>
        <li>JWT + bcrypt (authentication)</li>
        <li>Docker Compose (5 services)</li>
      </ul>
    </div>
  </div>

  <h2>Project Structure</h2>
  <p>The project contains two frontend applications:</p>
  <ul>
    <li><strong>frontend/</strong> — Harry Potter themed (9 pages: Landing, Dashboard, AnomalyDetail, MischiefList, Vault, Pensieve, OwlPost, Profile, Admin)</li>
    <li><strong>stitch-frontend/</strong> — OmniLedger tech/industrial themed (18+ pages including Fleet Management, Access Control, Global Feed, Integrations, Performance)</li>
  </ul>
  <p>Both frontends share the same <strong>backend/</strong> FastAPI application, ML engine, and database layer.</p>
</div>

<div class="page">
  <h1>2. Authentication & User Onboarding Workflow</h1>

  <div class="flow-row">
    <div class="flow-step">
      <div class="flow-num" style="background: #3B82F6;">1</div>
      <h4>Visit Site</h4>
      <p style="font-size:9pt">User navigates to the application URL</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #2D6A4F;">2</div>
      <h4>Login / Signup</h4>
      <p style="font-size:9pt">Register with name, email, password</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #D97706;">3</div>
      <h4>JWT Issued</h4>
      <p style="font-size:9pt">Backend returns JWT token (24h expiry)</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #7C3AED;">4</div>
      <h4>Onboarding Tour</h4>
      <p style="font-size:9pt">First-time users get a 5-step guided tour</p>
    </div>
  </div>

  <h2>Authentication Flow Details</h2>
  <div class="two-col">
    <div>
      <h4>Registration</h4>
      <ul>
        <li>POST /api/auth/register with name, email, password</li>
        <li>Password hashed with bcrypt (hash_password)</li>
        <li>User stored in SQLite users table</li>
        <li>Returns JWT + user info on success</li>
      </ul>
    </div>
    <div>
      <h4>Login</h4>
      <ul>
        <li>POST /api/auth/login with email, password</li>
        <li>Verifies against bcrypt hash</li>
        <li>Returns JWT with sub (user_id), email, name</li>
        <li>Token expires in 24 hours</li>
      </ul>
    </div>
  </div>

  <h4>Session Management</h4>
  <ul>
    <li>JWT stored in localStorage key 'marauders_token'</li>
    <li>User object stored in localStorage key 'marauders_session'</li>
    <li>Axios interceptor adds 'Authorization: Bearer <token>' to all requests</li>
    <li>ProtectedRoute component redirects to /login if not authenticated</li>
    <li>401 response from any API auto-clears session and redirects</li>
    <li>Rate limiting: 30 requests per 60s per IP on auth endpoints</li>
    <li>Backend verifies via FastAPI get_current_user dependency (HTTPBearer)</li>
    <li>User ownership checked on all endpoints (403 if mismatch)</li>
  </ul>
</div>

<div class="page">
  <h1>3. CSV Upload & Processing Workflow</h1>

  <div class="flow-row">
    <div class="flow-step">
      <div class="flow-num" style="background: #3B82F6;">1</div>
      <h4>Upload CSV</h4>
      <p style="font-size:9pt">Drag & drop CSV onto UploadZone component</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #2D6A4F;">2</div>
      <h4>Parse & Validate</h4>
      <p style="font-size:9pt">POST /api/upload - backend parses CSV</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #D97706;">3</div>
      <h4>Celery Task</h4>
      <p style="font-size:9pt">process_upload dispatched via Redis</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #7C3AED;">4</div>
      <h4>Feature Engineering</h4>
      <p style="font-size:9pt">48 features computed per transaction</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #DC2626;">5</div>
      <h4>ML Inference</h4>
      <p style="font-size:9pt">10-model ensemble runs detection</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #059669;">6</div>
      <h4>Results Stored</h4>
      <p style="font-size:9pt">Anomalies in VectorAI/SQLite</p>
    </div>
  </div>

  <h2>Upload Details</h2>
  <div class="two-col">
    <div>
      <h4>CSV Requirements</h4>
      <ul>
        <li>Required columns: amount, category, merchant, hour, day</li>
        <li>Uploaded via react-dropzone on the Landing page</li>
        <li>File sent to backend via FormData</li>
        <li>Backend parses with pandas (pd.read_csv)</li>
      </ul>
    </div>
    <div>
      <h4>Processing Pipeline</h4>
      <ul>
        <li>Celery task: process_upload → parse → features → detect</li>
        <li>Frontend polls GET /api/batches/:id/progress every 1s</li>
        <li>Polling timeout: 60 seconds</li>
        <li>On completion → navigate to /dashboard with toast</li>
        <li>Synchronous fallback if Celery/Redis unavailable</li>
      </ul>
    </div>
  </div>
</div>

<div class="page">
  <h1>4. ML Anomaly Detection Engine (Deep Dive)</h1>

  <h2>Model Architecture</h2>
  <table>
    <tr><th>Model</th><th>Type</th><th>Weight / Role</th></tr>
    <tr><td>Random Forest</td><td>Supervised Ensemble</td><td>0.15 — Core ensemble member</td></tr>
    <tr><td>Gradient Boosting</td><td>Supervised Ensemble</td><td>0.15 — Core ensemble member</td></tr>
    <tr><td>XGBoost</td><td>Gradient Boosting</td><td>0.25 — High-performance learner</td></tr>
    <tr><td>LightGBM</td><td>Gradient Boosting</td><td>0.25 — High-performance learner</td></tr>
    <tr><td>CatBoost</td><td>Gradient Boosting</td><td>0.10 — Supplementary</td></tr>
    <tr><td>Extra Trees</td><td>Supervised Ensemble</td><td>0.10 — Supplementary</td></tr>
    <tr><td>Isolation Forest</td><td>Unsupervised</td><td>Feature — Outlier score</td></tr>
    <tr><td>Local Outlier Factor</td><td>Unsupervised</td><td>Feature — Density score</td></tr>
    <tr><td>One-Class SVM</td><td>Unsupervised</td><td>Feature — Boundary score</td></tr>
    <tr><td>Rule-Based Scoring</td><td>Heuristic</td><td>0.25 weight — Domain rules</td></tr>
  </table>

  <h2>48 Engineered Features</h2>
  <div class="two-col">
    <div>
      <h4>Base Features (35)</h4>
      <ul>
        <li>Amount: raw, log-transformed, Z-score</li>
        <li>Category: one-hot encoded</li>
        <li>Merchant: encoded, new merchant flag</li>
        <li>Hour: raw, cyclical sin/cos encoding</li>
        <li>Day: raw, cyclical sin/cos, weekend flag</li>
        <li>Rolling: mean, std, count (7-day windows)</li>
        <li>Interaction: category×amount, hour×day</li>
        <li>Velocity: transaction count in time window</li>
      </ul>
    </div>
    <div>
      <h4>ML & Rule Scores (13)</h4>
      <ul>
        <li>Isolation Forest anomaly score</li>
        <li>Local Outlier Factor score</li>
        <li>One-Class SVM decision function</li>
        <li>Rule Score: heuristic anomaly indicator</li>
        <li>Amount threshold violation score</li>
        <li>Category spending spike score</li>
        <li>Unusual hour indicator</li>
        <li>New merchant detection score</li>
        <li>Rolling average deviation score</li>
      </ul>
    </div>
  </div>

  <h2>Scoring Formula</h2>
  <p><strong>final_score = 0.75 x ensemble_score + 0.25 x rule_score</strong></p>
  <ul>
    <li>ensemble_score = weighted average of 9 ML model predictions</li>
    <li>rule_score = heuristic-based anomaly indicators (0-1)</li>
    <li>Threshold: final_score > 0.5 → flagged as anomaly</li>
    <li>Severity: <span class="badge badge-green">Peeves (low)</span> <span class="badge badge-amber">Boggart (medium)</span> <span class="badge badge-red">Dementor (high)</span></li>
  </ul>
  <p>Feature scaling via StandardScaler. All models serialized with joblib.</p>
</div>

<div class="page">
  <h1>5. Dashboard & Marauder's Map Workflow</h1>

  <p>After CSV processing completes, the user is redirected to the Dashboard where results are visualized on the interactive Marauder's Map.</p>

  <h2>Map Locations (Spending Categories)</h2>
  <table>
    <tr><th>Location</th><th>Category</th><th>Color</th></tr>
    <tr><td>Hogwarts</td><td>Food & Dining</td><td>Purple</td></tr>
    <tr><td>Hogsmeade</td><td>Shopping</td><td>Blue</td></tr>
    <tr><td>Gringotts</td><td>Bills & Utilities</td><td>Gold</td></tr>
    <tr><td>Diagon Alley</td><td>Entertainment</td><td>Teal/Green</td></tr>
    <tr><td>Platform 9 3/4</td><td>Travel</td><td>Red</td></tr>
  </table>

  <h2>Interactive Features</h2>
  <div class="two-col">
    <div>
      <ul>
        <li><strong>Animated Footprints</strong> — Transactions appear as walking footprints traversing the map between locations</li>
        <li><strong>Anomaly Pulsing</strong> — Detected anomalies glow red and pulse to draw immediate attention</li>
        <li><strong>Severity Colors</strong> — Peeves (green/low), Boggart (amber/medium), Dementor (red/high)</li>
        <li><strong>Click to Inspect</strong> — Click any footprint or anomaly to navigate to the detail page</li>
      </ul>
    </div>
    <div>
      <ul>
        <li><strong>Category Filter Tabs</strong> — Moony (Food), Wormtail (Shopping), Padfoot (Bills), Prongs (All)</li>
        <li><strong>Spending Trend Chart</strong> — Bottom panel shows Recharts line chart of spending over time</li>
        <li><strong>Stats Bar</strong> — Top bar shows total anomalies, high/medium/low counts, total financial impact</li>
        <li><strong>Right Sidebar</strong> — Anomaly cards list with severity badges</li>
      </ul>
    </div>
  </div>

  <p>The MaraudersMap component (257 lines) is an SVG-based interactive visualization built with React and Framer Motion for animations. The SpendChart component uses Recharts for the spending trend line chart.</p>
</div>

<div class="page">
  <h1>6. Anomaly Investigation Workflow</h1>

  <div class="flow-row">
    <div class="flow-step">
      <div class="flow-num" style="background: #3B82F6;">1</div>
      <h4>Click Anomaly</h4>
      <p style="font-size:9pt">From map or card → /anomaly/:id</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #2D6A4F;">2</div>
      <h4>Score Gauges</h4>
      <p style="font-size:9pt">ML Model + Rule + Final scores</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #7C3AED;">3</div>
      <h4>AI Narrative</h4>
      <p style="font-size:9pt">Gemini generates HP-themed story</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #D97706;">4</div>
      <h4>Audio Narration</h4>
      <p style="font-size:9pt">ElevenLabs TTS playback</p>
    </div>
    <div class="flow-arrow">→</div>
    <div class="flow-step">
      <div class="flow-num" style="background: #DC2626;">5</div>
      <h4>Review & Act</h4>
      <p style="font-size:9pt">Mark valid/mischief, export, share</p>
    </div>
  </div>

  <h2>Anomaly Detail Page Components</h2>
  <div class="two-col">
    <div>
      <h4>Investigation Tools</h4>
      <ul>
        <li><strong>ScoreGauge (x3)</strong> — Circular gauges showing ML Model Score, Rule Score, and Final Mischief Score (0-100%)</li>
        <li><strong>SeverityBadge</strong> — Color-coded severity indicator (Dementor/Boggart/Peeves)</li>
        <li><strong>NarrativeCard</strong> — Gemini-generated AI narrative with typewriter animation effect</li>
        <li><strong>AudioPlayer</strong> — Waveform audio player for TTS narration (play/pause controls)</li>
      </ul>
    </div>
    <div>
      <h4>Actions & Context</h4>
      <ul>
        <li><strong>Status Management</strong> — Toggle between 'Valid' (confirmed legitimate) and 'Mischief' (confirmed fraud)</li>
        <li><strong>Triggered Rules</strong> — Tags showing which rules the anomaly violated</li>
        <li><strong>Related Transactions</strong> — Nearby transactions for context and comparison</li>
        <li><strong>Export & Share</strong> — Export HTML report via pdfExport.ts, share link via clipboard</li>
      </ul>
    </div>
  </div>

  <h2>Narrative Generation Flow</h2>
  <ol>
    <li>Frontend calls GET /api/narratives/:id</li>
    <li>Backend checks cache (existing narrative in DB)</li>
    <li>If not cached: calls Google Gemini 2.0 Flash API with anomaly data</li>
    <li>Gemini generates Harry Potter-themed narrative explaining the anomaly</li>
    <li>Narrative stored in database for future requests</li>
    <li>Frontend displays with typewriter animation</li>
    <li>Graceful fallback: static template if Gemini API unavailable</li>
  </ol>

  <h2>TTS Audio Flow</h2>
  <ol>
    <li>User clicks play on AudioPlayer widget</li>
    <li>Frontend calls GET /api/narratives/:id/audio</li>
    <li>Backend checks cache for existing MP3 audio</li>
    <li>If not cached: calls ElevenLabs TTS API (voice: Roger)</li>
    <li>Audio streamed back as MP3 (Content-Type: audio/mpeg)</li>
    <li>Sequential request locking per anomaly ID prevents duplicates</li>
    <li>Graceful fallback: narrative displayed as text only</li>
  </ol>
</div>

<div class="page">
  <h1>7. Voice Chat & Power User Features</h1>

  <div class="section-header">Voice Chat Assistant</div>
  <div class="two-col">
    <div>
      <ul>
        <li>Persistent FAB (floating action button) in bottom-right corner of every page</li>
        <li>Supports voice input via Web Speech API (SpeechRecognition)</li>
        <li>Text input fallback for environments without microphone access</li>
        <li>Sends POST /api/chat/message with page context (anomaly_id, batch_id)</li>
      </ul>
    </div>
    <div>
      <ul>
        <li>Gemini responds in Marauder's Map persona character</li>
        <li>ElevenLabs streams audio response back (X-Chat-Response header + audio/mpeg)</li>
        <li>Chat history preserved during session</li>
        <li>Example queries: "What anomalies were detected?", "Explain this transaction", "Summarize my spending"</li>
      </ul>
    </div>
  </div>

  <div class="section-header">Command Palette (Cmd+K / Ctrl+K)</div>
  <ul>
    <li>Quick navigation to any page in the application via fuzzy search</li>
    <li>Search transactions, anomalies, and investigations</li>
    <li>Toggle dark mode directly from command palette</li>
    <li>Access keyboard shortcuts reference modal</li>
    <li>Context-aware commands based on the current page</li>
  </ul>

  <div class="section-header">Dark Mode & Theming</div>
  <ul>
    <li>Full dark/light theme toggle via DarkModeToggle component</li>
    <li>Persistent theme preference stored in localStorage ('theme' key)</li>
    <li>Seamless transitions via Tailwind CSS dark mode classes</li>
    <li>Parchment aging effect in the Harry Potter theme</li>
    <li>Consistent design across all 18+ pages</li>
  </ul>

  <div class="section-header">PDF Export & Sharing</div>
  <ul>
    <li>Export anomaly reports as formatted HTML documents</li>
    <li>Share anomaly details via clipboard link</li>
    <li>PDF includes: score gauges, AI narrative, triggered rules, transaction details</li>
    <li>Export function in AnomalyDetail page via pdfExport.ts service</li>
  </ul>
</div>

<div class="page">
  <h1>8. Supporting Pages Workflow</h1>

  <table>
    <tr><th>Page</th><th>Route</th><th>Purpose</th><th>Key Features</th></tr>
    <tr>
      <td>MischiefList</td>
      <td>/ledger</td>
      <td>Historical anomaly table</td>
      <td>Search, filter by severity, sort by date, paginated list</td>
    </tr>
    <tr>
      <td>Vault</td>
      <td>/vault</td>
      <td>Gringotts-style spending overview</td>
      <td>Total balance, category breakdown with ProgressRing, recent anomalies</td>
    </tr>
    <tr>
      <td>Pensieve</td>
      <td>/pensieve</td>
      <td>Deep spending analysis</td>
      <td>Spending trend line charts, time range selectors, risk breakdown pie chart</td>
    </tr>
    <tr>
      <td>Owl Post</td>
      <td>/owl-post</td>
      <td>Notification center</td>
      <td>Read/unread anomaly alerts, categorized by severity with timestamps</td>
    </tr>
    <tr>
      <td>Profile</td>
      <td>/profile</td>
      <td>Wizard's dossier</td>
      <td>User stats, spending by category rings, recent cases list</td>
    </tr>
    <tr>
      <td>Admin</td>
      <td>/admin</td>
      <td>System settings</td>
      <td>Theme toggle, parchment aging, mischief taxonomy editor, API status dashboard</td>
    </tr>
  </table>

  <h2>Navigation</h2>
  <ul>
    <li>Persistent SidebarNav component on all pages</li>
    <li>Header with UserMenu dropdown (profile, settings, logout)</li>
    <li>Command Palette (Cmd+K) for quick page navigation</li>
    <li>Keyboard shortcut help modal</li>
    <li>PageTransition animations via Framer Motion AnimatePresence</li>
  </ul>

  <h2>State Management</h2>
  <ul>
    <li><strong>React Context:</strong> AuthContext (user/JWT), AppContext (batchId, userId), ToastContext (notifications)</li>
    <li><strong>TanStack Query:</strong> All API data managed via useQuery + useMutation with automatic caching and invalidation</li>
    <li><strong>Query Keys:</strong> ['anomalies'], ['narrative', id], ['audio', id], ['spending-day'], ['spending-category'], ['batches'], ['transactions']</li>
    <li><strong>localStorage:</strong> marauders_token, marauders_session, marauders_onboarded, theme</li>
  </ul>
</div>

<div class="page">
  <h1>9. End-to-End Data Flow Architecture</h1>

  <table>
    <tr><th>Layer</th><th>Technology</th><th>Responsibilities</th></tr>
    <tr>
      <td>User Browser</td>
      <td>React 19 + TypeScript</td>
      <td>BrowserRouter → AuthProvider → QueryClientProvider → AppProvider. TanStack Query for all API data. Components: SidebarNav, UploadZone, MaraudersMap, VoiceChatWidget</td>
    </tr>
    <tr>
      <td>API Gateway</td>
      <td>FastAPI + Uvicorn</td>
      <td>17+ REST endpoints. JWT verification via get_current_user. CORS middleware, rate limiting (30 req/60s per IP). Async handlers with HTTPX for external API calls.</td>
    </tr>
    <tr>
      <td>Background Processing</td>
      <td>Celery + Redis</td>
      <td>process_upload: parse CSV → feature engineering → ML inference. Frontend polls progress every 1s. Synchronous fallback if unavailable.</td>
    </tr>
    <tr>
      <td>ML Engine</td>
      <td>scikit-learn + XGBoost + LightGBM + CatBoost</td>
      <td>10 models: 6 supervised + 3 unsupervised + 1 rule engine. 48 features per transaction → scaled → ensemble prediction. Hybrid score: 0.75*ML + 0.25*rules.</td>
    </tr>
    <tr>
      <td>Data Layer</td>
      <td>Actian VectorAI + SQLite</td>
      <td>VectorAI: primary vector database with hybrid search. SQLite: automatic fallback + user auth storage. Auto-detection of server capabilities.</td>
    </tr>
    <tr>
      <td>External AI</td>
      <td>Gemini 2.0 Flash + ElevenLabs</td>
      <td>Gemini: narrative generation + voice chat responses. ElevenLabs: TTS narration streaming MP3. Thread-safe locking per anomaly. Graceful degradation.</td>
    </tr>
  </table>

  <h2>Deployment Architecture (Docker Compose)</h2>
  <table>
    <tr><th>Service</th><th>Role</th><th>Port</th></tr>
    <tr><td>Frontend</td><td>React 19 + Nginx static serving</td><td>3000</td></tr>
    <tr><td>Backend</td><td>FastAPI + Uvicorn ASGI server</td><td>8000</td></tr>
    <tr><td>Redis</td><td>Celery message broker + cache</td><td>6379</td></tr>
    <tr><td>Celery Worker</td><td>Background task processing</td><td>—</td></tr>
    <tr><td>VectorAI</td><td>Vector database (Actian community)</td><td>8888</td></tr>
  </table>
  <p>Start command: <code>docker-compose up --build</code></p>
</div>

<div class="page">
  <h1>10. Performance Metrics & Model Validation</h1>

  <h2>Key Metrics</h2>
  <div class="grid">
    <div class="grid-item">
      <h4>F1 Score</h4>
      <p style="font-size:24pt;color:#2D6A4F;font-weight:bold;margin:0;">0.873</p>
      <p style="font-size:9pt;">Test Set</p>
    </div>
    <div class="grid-item">
      <h4>Precision</h4>
      <p style="font-size:24pt;color:#3B82F6;font-weight:bold;margin:0;">0.977</p>
      <p style="font-size:9pt;">Test Set</p>
    </div>
    <div class="grid-item">
      <h4>Recall</h4>
      <p style="font-size:24pt;color:#D97706;font-weight:bold;margin:0;">0.897</p>
      <p style="font-size:9pt;">Test Set</p>
    </div>
    <div class="grid-item">
      <h4>CV F1</h4>
      <p style="font-size:24pt;color:#7C3AED;font-weight:bold;margin:0;">0.8808</p>
      <p style="font-size:9pt;">Mean +/- 0.0017</p>
    </div>
  </div>

  <h2>Ensemble Model Weights</h2>
  <table>
    <tr><th>Model</th><th>Weight</th></tr>
    <tr><td>XGBoost</td><td>0.25</td></tr>
    <tr><td>LightGBM</td><td>0.25</td></tr>
    <tr><td>Random Forest</td><td>0.15</td></tr>
    <tr><td>Gradient Boosting</td><td>0.15</td></tr>
    <tr><td>CatBoost</td><td>0.10</td></tr>
    <tr><td>Extra Trees</td><td>0.10</td></tr>
    <tr><td>Rule-Based Scoring</td><td>0.25 (of hybrid formula)</td></tr>
  </table>

  <h2>Feature Summary</h2>
  <ul>
    <li><strong>35 Base Features:</strong> amount, category, merchant, hour, day, cyclical time encoding (sin/cos), log transforms, Z-scores, rolling statistics (mean, std, count), interaction features (category×amount, hour×day)</li>
    <li><strong>3 Unsupervised Scores:</strong> Isolation Forest score, LOF score, OCSVM score</li>
    <li><strong>1 Rule Score:</strong> heuristic-based anomaly indicators</li>
    <li><strong>Feature scaling:</strong> StandardScaler normalization</li>
  </ul>
</div>

<div class="page">
  <h1>11. Complete Workflow Summary</h1>

  <h2>End-to-End User Journey</h2>
  <ol>
    <li><strong>Authentication</strong> — User signs up or logs in. JWT stored in localStorage. Protected routes guard all pages.</li>
    <li><strong>CSV Upload</strong> — User drag-and-drops a CSV file onto the Landing page's UploadZone. File sent to backend via POST /api/upload. Celery worker processes asynchronously with progress polling.</li>
    <li><strong>ML Detection</strong> — 48 features engineered per transaction. 10-model ensemble runs anomaly detection. Hybrid scoring combines ML predictions with rule-based heuristics.</li>
    <li><strong>Map Display</strong> — Results shown on interactive Marauder's Map SVG. Footprints animate, anomalies pulse with severity colors, category filter tabs available.</li>
    <li><strong>Investigation</strong> — User clicks anomaly for detail view. Three score gauges, Gemini-generated AI narrative with typewriter effect, ElevenLabs TTS audio narration.</li>
    <li><strong>Actions</strong> — Mark as 'Valid' or 'Mischief'. Export HTML report. Share via clipboard link. Use voice chat for natural language queries.</li>
  </ol>

  <h2>Data Flow Summary</h2>
  <p style="font-family:'Courier New',monospace;font-size:9pt;background:#2C1810;color:#D4AF37;padding:10px;border-radius:4px;">
    User → React Frontend → FastAPI Backend → Redis/Celery → ML Ensemble → VectorAI/SQLite → Gemini AI → ElevenLabs TTS → User
  </p>

  <h2>Key Differentiators</h2>
  <div class="grid">
    <div class="grid-item">
      <h4>10-Model Ensemble</h4>
      <p>Combines 6 supervised + 3 unsupervised + rule engine for robust detection</p>
    </div>
    <div class="grid-item">
      <h4>48 Features</h4>
      <p>Deep feature engineering including rolling stats, cyclical encoding, and interaction features</p>
    </div>
    <div class="grid-item">
      <h4>AI Narratives</h4>
      <p>Gemini generates contextual, engaging explanations for every anomaly</p>
    </div>
    <div class="grid-item">
      <h4>Interactive Map</h4>
      <p>SVG-based Marauder's Map with animated footprints and real-time updates</p>
    </div>
    <div class="grid-item">
      <h4>Voice Interface</h4>
      <p>Voice chat + TTS narration for accessibility and natural interaction</p>
    </div>
    <div class="grid-item">
      <h4>Dual Frontend</h4>
      <p>Two complete frontends (HP theme + OmniLedger tech theme) sharing one backend</p>
    </div>
  </div>

  <div style="text-align:center;margin-top:30px;padding:15px;background:#1A0F0A;color:#D4AF37;border-radius:4px;">
    <p style="font-size:14pt;font-style:italic;margin:0;">"Mischief Managed."</p>
    <p style="font-size:10pt;margin:5px 0 0 0;">The Marauder's Ledger v1.0 — AI-Powered Financial Anomaly Detection</p>
  </div>
</div>

</body>
</html>
"""

pdf_path = "/home/saikat/Marauders-Ledger/Marauders_Ledger_Workflow_Report.pdf"
HTML(string=pdf_html).write_pdf(pdf_path)
print(f"PDF saved: {pdf_path}")

print("\n=== DONE ===")
print(f"PPTX: {pptx_path}")
print(f"PDF:  {pdf_path}")
